from __future__ import annotations

import hashlib
import io
import json
from decimal import Decimal
from typing import Any

from django.http import HttpResponse
from django.utils import timezone
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from .models import (
    FileRecord,
    ProposalCloseout,
    ProposalFinding,
    ProposalPlan,
    ProposalRequirement,
    ProposalReview,
    ProposalSubmissionSnapshot,
)
from .proposal_execution import proposal_execution_payload


def _user_name(user) -> str:
    if not user:
        return ""
    full = (getattr(user, "get_full_name", lambda: "")() or "").strip()
    return full or getattr(user, "email", "") or getattr(user, "username", "") or str(user.pk)


def _opportunity_snapshot(opportunity) -> dict[str, Any]:
    return {
        "source_id": opportunity.source_id,
        "title": opportunity.title,
        "agency": opportunity.agency,
        "solicitation_number": opportunity.solicitation_number,
        "naics_code": opportunity.naics_code,
        "psc_code": opportunity.psc_code,
        "response_deadline": opportunity.response_deadline.isoformat() if opportunity.response_deadline else None,
        "updated_at": opportunity.updated_at.isoformat() if opportunity.updated_at else None,
        "source_url": opportunity.source_url,
    }


def _requirement_snapshot(plan: ProposalPlan) -> list[dict[str, Any]]:
    return [
        {
            "id": row.id,
            "key": row.key,
            "requirement": row.requirement,
            "source": row.source,
            "source_kind": row.source_kind,
            "evidence": row.evidence,
            "status": row.status,
            "owner": _user_name(row.owner),
            "due_at": row.due_at.isoformat() if row.due_at else None,
            "notes": row.notes,
        }
        for row in plan.requirements.select_related("owner").all()
    ]


def _review_snapshot(plan: ProposalPlan) -> list[dict[str, Any]]:
    return [
        {
            "id": row.id,
            "review_type": row.review_type,
            "label": row.get_review_type_display(),
            "status": row.status,
            "owner": _user_name(row.owner),
            "target_at": row.target_at.isoformat() if row.target_at else None,
            "completed_at": row.completed_at.isoformat() if row.completed_at else None,
            "summary": row.summary,
        }
        for row in plan.reviews.select_related("owner").all()
    ]


def _finding_snapshot(plan: ProposalPlan) -> list[dict[str, Any]]:
    return [
        {
            "id": row.id,
            "severity": row.severity,
            "title": row.title,
            "detail": row.detail,
            "status": row.status,
            "owner": _user_name(row.owner),
            "due_at": row.due_at.isoformat() if row.due_at else None,
            "review_id": row.review_id,
            "requirement_id": row.requirement_id,
        }
        for row in plan.findings.select_related("owner").all()
    ]


def _file_manifest(plan: ProposalPlan) -> list[dict[str, Any]]:
    rows = FileRecord.objects.filter(
        organization=plan.organization,
        opportunity=plan.opportunity,
    ).order_by("name", "-version")
    return [
        {
            "id": row.id,
            "name": row.name,
            "file_type": row.file_type,
            "version": row.version,
            "size_bytes": row.size_bytes,
            "checksum": row.checksum,
            "source": row.source,
            "source_url": row.source_url,
        }
        for row in rows
    ]


def _hash_payload(payload: dict[str, Any]) -> str:
    raw = json.dumps(payload, sort_keys=True, separators=(",", ":"), default=str).encode("utf-8")
    return hashlib.sha256(raw).hexdigest()


def submission_readiness(plan: ProposalPlan, execution: dict[str, Any]) -> dict[str, Any]:
    blockers: list[str] = []
    if not execution["plan"]["submission_ready"]:
        blockers.append("Proposal Execution is not yet submission ready.")
    if execution["amendment_impact"]["changed"]:
        blockers.append("The solicitation changed after the last reviewed amendment baseline.")
    if execution["counts"]["requirements_closed"] < execution["counts"]["requirements_total"]:
        blockers.append("One or more proposal requirements still need human disposition.")
    if execution["counts"]["reviews_passed"] < execution["counts"]["reviews_total"]:
        blockers.append("All Pink/Red/Gold/Final review gates have not passed.")
    if execution["counts"]["open_findings"]:
        blockers.append("Open proposal review findings remain.")
    if not plan.final_submission_verified:
        blockers.append("Final submission verification has not been completed by a human user.")
    files = _file_manifest(plan)
    if not files:
        blockers.append("No final proposal files are registered for this opportunity.")
    return {
        "ready": not blockers,
        "blockers": blockers,
        "file_count": len(files),
        "deadline": plan.opportunity.response_deadline,
    }


def build_submission_control(*, organization, opportunity, user) -> dict[str, Any]:
    from .proposal_execution import ensure_proposal_execution
    plan = ensure_proposal_execution(organization=organization, opportunity=opportunity, user=user)
    execution = proposal_execution_payload(organization=organization, opportunity=opportunity, user=user)
    readiness = submission_readiness(plan, execution)
    closeout, _ = ProposalCloseout.objects.get_or_create(plan=plan, defaults={"updated_by": user})
    snapshots = [
        {
            "id": row.id,
            "sequence": row.sequence,
            "submitted_at": row.submitted_at,
            "submitted_by": _user_name(row.submitted_by),
            "delivery_method": row.delivery_method,
            "confirmation_reference": row.confirmation_reference,
            "snapshot_hash": row.snapshot_hash,
            "file_count": len(row.file_manifest or []),
        }
        for row in plan.submission_snapshots.select_related("submitted_by").all()[:25]
    ]
    return {
        "plan": {
            "id": plan.id,
            "status": plan.status,
            "submission_method": plan.submission_method,
            "final_submission_verified": plan.final_submission_verified,
        },
        "submission_readiness": readiness,
        "snapshots": snapshots,
        "closeout": {
            "id": closeout.id,
            "status": closeout.status,
            "awardee": closeout.awardee,
            "award_value": closeout.award_value,
            "award_date": closeout.award_date,
            "debrief_requested": closeout.debrief_requested,
            "debrief_received": closeout.debrief_received,
            "win_loss_reason": closeout.win_loss_reason,
            "customer_feedback": closeout.customer_feedback,
            "strengths": closeout.strengths,
            "weaknesses": closeout.weaknesses,
            "lessons_learned": closeout.lessons_learned,
        },
        "exports": [
            {"format": "pdf", "label": "Executive Opportunity Brief PDF"},
            {"format": "xlsx", "label": "Compliance Matrix XLSX"},
            {"format": "pptx", "label": "Management Summary PowerPoint"},
        ],
        "warning": "ForgeGov preserves submission snapshots as historical records. Verify the official portal receipt and agency instructions independently.",
    }


def create_submission_snapshot(*, plan: ProposalPlan, user, confirmation_reference: str = "", notes: str = "") -> ProposalSubmissionSnapshot:
    execution = proposal_execution_payload(organization=plan.organization, opportunity=plan.opportunity, user=user)
    readiness = submission_readiness(plan, execution)
    if not readiness["ready"]:
        raise ValueError("Submission snapshot cannot be created until all submission-readiness blockers are cleared.")

    payload = {
        "opportunity": _opportunity_snapshot(plan.opportunity),
        "requirements": _requirement_snapshot(plan),
        "reviews": _review_snapshot(plan),
        "findings": _finding_snapshot(plan),
        "files": _file_manifest(plan),
        "amendment": plan.amendment_baseline or {},
        "delivery_method": plan.submission_method,
        "confirmation_reference": confirmation_reference.strip(),
        "notes": notes.strip(),
    }
    sequence = (plan.submission_snapshots.order_by("-sequence").values_list("sequence", flat=True).first() or 0) + 1
    snapshot = ProposalSubmissionSnapshot.objects.create(
        plan=plan,
        sequence=sequence,
        submitted_at=timezone.now(),
        submitted_by=user,
        delivery_method=plan.submission_method,
        confirmation_reference=confirmation_reference.strip()[:500],
        notes=notes,
        opportunity_snapshot=payload["opportunity"],
        requirement_snapshot=payload["requirements"],
        review_snapshot=payload["reviews"],
        finding_snapshot=payload["findings"],
        file_manifest=payload["files"],
        amendment_snapshot=payload["amendment"],
        snapshot_hash=_hash_payload(payload),
    )
    plan.status = ProposalPlan.Status.SUBMITTED
    plan.save(update_fields=["status", "updated_at"])
    ProposalCloseout.objects.update_or_create(
        plan=plan,
        defaults={"status": ProposalCloseout.Status.SUBMITTED, "updated_by": user},
    )
    return snapshot


def update_closeout(closeout: ProposalCloseout, payload: dict[str, Any], user) -> None:
    allowed = {
        "status", "awardee", "award_value", "award_date", "debrief_requested",
        "debrief_received", "win_loss_reason", "customer_feedback",
        "strengths", "weaknesses", "lessons_learned",
    }
    for key in allowed:
        if key not in payload:
            continue
        value = payload[key]
        if key == "status" and value not in {choice for choice, _ in ProposalCloseout.Status.choices}:
            raise ValueError("Invalid closeout status.")
        if key == "award_value" and value in ("", None):
            value = None
        if key in {"strengths", "weaknesses", "lessons_learned"} and not isinstance(value, list):
            value = [str(value)] if value else []
        setattr(closeout, key, value)
    closeout.updated_by = user
    closeout.save()


def export_submission_control(*, plan: ProposalPlan, fmt: str) -> HttpResponse:
    fmt = fmt.lower()
    opportunity = plan.opportunity
    requirements = _requirement_snapshot(plan)
    reviews = _review_snapshot(plan)

    if fmt == "xlsx":
        wb = Workbook()
        ws = wb.active
        ws.title = "Compliance Matrix"
        ws.append(["Requirement", "Status", "Owner", "Source", "Evidence", "Due"])
        for row in requirements:
            ws.append([row["requirement"], row["status"], row["owner"], row["source"], row["evidence"], row["due_at"] or ""])
        rev = wb.create_sheet("Reviews")
        rev.append(["Review", "Status", "Owner", "Target", "Completed", "Summary"])
        for row in reviews:
            rev.append([row["label"], row["status"], row["owner"], row["target_at"] or "", row["completed_at"] or "", row["summary"]])
        output = io.BytesIO()
        wb.save(output)
        response = HttpResponse(output.getvalue(), content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        response["Content-Disposition"] = 'attachment; filename="forgegov-compliance-matrix.xlsx"'
        return response

    if fmt == "pptx":
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = opportunity.title or "ForgeGov Opportunity"
        slide.placeholders[1].text = f"{opportunity.agency or 'Agency'}\n{opportunity.solicitation_number or opportunity.source_id}"
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Proposal Status"
        body = slide.placeholders[1].text_frame
        body.text = f"Plan status: {plan.status}"
        for line in [
            f"Submission method: {plan.submission_method or 'Not specified'}",
            f"Requirements: {len(requirements)}",
            f"Reviews: {len(reviews)}",
            f"Final verification: {'Complete' if plan.final_submission_verified else 'Pending'}",
        ]:
            p = body.add_paragraph()
            p.text = line
        output = io.BytesIO()
        prs.save(output)
        response = HttpResponse(output.getvalue(), content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        response["Content-Disposition"] = 'attachment; filename="forgegov-management-summary.pptx"'
        return response

    if fmt == "pdf":
        output = io.BytesIO()
        c = canvas.Canvas(output, pagesize=letter)
        width, height = letter
        y = height - 54
        c.setFont("Helvetica-Bold", 16)
        c.drawString(54, y, "ForgeGov Executive Opportunity Brief")
        y -= 28
        c.setFont("Helvetica-Bold", 11)
        for label, value in [
            ("Opportunity", opportunity.title),
            ("Agency", opportunity.agency),
            ("Solicitation", opportunity.solicitation_number or opportunity.source_id),
            ("Proposal status", plan.status),
            ("Submission method", plan.submission_method or "Not specified"),
        ]:
            c.drawString(54, y, f"{label}:")
            c.setFont("Helvetica", 10)
            c.drawString(150, y, str(value or "—")[:85])
            c.setFont("Helvetica-Bold", 11)
            y -= 18
        y -= 10
        c.drawString(54, y, "Compliance Summary")
        y -= 18
        c.setFont("Helvetica", 9)
        for row in requirements[:22]:
            line = f"[{row['status']}] {row['requirement']}"
            c.drawString(60, y, line[:105])
            y -= 14
            if y < 60:
                c.showPage()
                y = height - 54
                c.setFont("Helvetica", 9)
        c.save()
        response = HttpResponse(output.getvalue(), content_type="application/pdf")
        response["Content-Disposition"] = 'attachment; filename="forgegov-executive-opportunity-brief.pdf"'
        return response

    raise ValueError("Unsupported export format.")
